import sys
import os
import anndata as ad
import pandas as pd
import scanpy as sc
import seaborn as sns; sns.set(color_codes=True)
import matplotlib
#matplotlib.rcParams["image.interpolation"] = None
import matplotlib.pyplot as plt
import matplotlib.pylab as pylab
from skimage.morphology import (erosion, dilation)
from skimage.morphology import disk  
footprint = disk(1)
from PIL import Image
Image.MAX_IMAGE_PIXELS = 3000000000
from skimage.measure import label, regionprops, regionprops_table
from tifffile import imread, imwrite
from tifffile import TiffFile
import numpy as np
import scipy.stats as stats
import time
import pickle
#import rapids_scanpy_funcs as rs
import subprocess as sb
from xml.etree import ElementTree
import gc
import math
from time import sleep
from tqdm import tqdm
import scipy.io
import numpy as np
from pptx import Presentation
import cudf
import cugraph
import rmm
import cupy
from PIL import Image
from skimage import morphology
import random
import hdf5storage

#cupy.cuda.set_allocator(rmm.rmm_cupy_allocator)


def Create_Marker_Info_Table(Projects_path, project_name, sample_ID, includemarker='True'):
    
    if not os.path.exists(os.path.join(Projects_path, project_name, 'Image-Data',sample_ID)):
           os.makedirs(os.path.join(Projects_path, project_name, 'Image-Data',sample_ID))

    mypath_img=os.path.join(Projects_path, project_name, 'Image-Data',sample_ID,sample_ID+'.qptiff')
    tif = TiffFile(mypath_img)

    marker_info_table = pd.DataFrame()
    marker_info_table = pd.DataFrame(columns=['marker_names', 'compartment', 'use_for_clustering', 'QC_passed'])

    FirstP_widthsize=tif.pages[00].imagewidth

    Biomarker_list=[]
    with TiffFile(mypath_img) as tif:
           for page in tif.series[0].pages:
                if page.imagewidth != FirstP_widthsize:
                        break
                Biomarkername= ElementTree.fromstring(page.description).find('Biomarker').text
                Biomarker_list.append(Biomarkername.replace(" ", ""))

    compartment=[]
    use_for_clustering=[]
    
    for i in range(np.size(Biomarker_list)):
        if (Biomarker_list[i].upper().startswith('DAPI')) or (Biomarker_list[i].upper()=='FOXP3') or (Biomarker_list[i].upper()=='PAX5') or (Biomarker_list[i].upper()=='KI67') or (Biomarker_list[i].upper()=='PCNA') or (Biomarker_list[i].upper()=='TP63'):
            compartment.append('nuclear')
        else:
            compartment.append('Membrane')
            
        if (Biomarker_list[i].upper().startswith('DAPI')) or (Biomarker_list[i].upper()=='HLA-A'):  
            use_for_clustering.append('False')
        else: 
            use_for_clustering.append(includemarker)

    marker_info_table['marker_names']= Biomarker_list
    marker_info_table['compartment']= compartment
    marker_info_table['use_for_clustering']= use_for_clustering
    marker_info_table['QC_passed']= 'Yes'
    marker_info_table['min_cutoff']= 0
    marker_info_table['max_cutoff']= 255
    marker_info_table.to_csv(os.path.join(Projects_path, project_name, 'Image-Data',sample_ID)+'/'+sample_ID+'_marker_info_table.csv', index = False)


def Stardist_segmentation(CodexObj, CLAHE=1, factor=1, tilesize=3096, min_area=5, runby='GPU', mask_already_exist=0):
    # Read the DAPI Image
    startTime = time.time()
    DAPI_im = imread(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["qptiff_path"],
                                  CodexObj["sample_ID"] + '.qptiff'), key=[0])
    
    # Create two folders, one called Dapi_im, and one called segmented_mask
    if not os.path.exists(os.path.join(
            os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"], 'Output-data'),
            'Dapi_im')):
        os.makedirs(os.path.join(
            os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"], 'Output-data'),
            'Dapi_im'))
    if not os.path.exists(os.path.join(
            os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"], 'Output-data'),
            'segmented_mask')):
        os.makedirs(os.path.join(
            os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"], 'Output-data'),
            'segmented_mask'))

    #write the Dapi image in the Dapi_im folder so that stardist executable can read it
    imwrite(os.path.join(
        os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"], 'Output-data',
                     'Dapi_im', CodexObj["sample_ID"] + '.tif')), DAPI_im)
    
    # Prepare the input and output path for stardist executable
    inputpath = os.path.join(
        os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"], 'Output-data',
                     'Dapi_im', CodexObj["sample_ID"] + '.tif'))
    outputpath = os.path.join(
        os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"], 'Output-data',
                     'segmented_mask', CodexObj["sample_ID"] + '_segmented.tif'))
    print(outputpath)
    if mask_already_exist ==0:
        print('Stardist segmentation starts.....')
        if runby=='GPU':
            os.system(f'./Run-NewStardist-v6-linux --input {inputpath} --output {outputpath} --CLAHE {CLAHE} --factor {factor} --tilesize {tilesize}')
        else:
            os.system(f'./Run-NewStardist-v6-linux-cpu --input {inputpath} --output {outputpath} --CLAHE {CLAHE} --factor {factor} --tilesize {tilesize}')
        print('Stardist segmentation ends.....')

    CodexObj["time_Stardist_segmentation"]= time.time() - startTime
    Dapi_mask= plt.imread(outputpath)

    # Resize if factor different than 1: This will be added to Stardist exectuable in future release  
    if factor !=1:
         footprint = disk(1)
         Labeled_Dapi_mask=label(Dapi_mask, connectivity=1)
         Labeled_Dapi_mask=Image.fromarray(np.asarray(Labeled_Dapi_mask).astype('float32'))
         Labeled_Dapi_maskA=np.array(Labeled_Dapi_mask.resize((DAPI_im.shape[1], DAPI_im.shape[0]), Image.NEAREST))
         if factor>1:    #it means the resulted mask from stardist is bigger and we need to make it smaller
            Labeled_Dapi_maskAD=dilation(Labeled_Dapi_maskA, footprint)
            Labeled_Dapi_mask_R=(Labeled_Dapi_maskAD-Labeled_Dapi_maskA)>0
            Labeled_Dapi_mask_R=np.invert(Labeled_Dapi_mask_R).astype(int)
            Labeled_Dapi_maskA=Labeled_Dapi_mask_R*Labeled_Dapi_maskA
         Dapi_mask = Labeled_Dapi_maskA.copy()
         dilated_labels = dilation(Labeled_Dapi_maskA, footprint)
         Boundary = dilated_labels - Labeled_Dapi_maskA
         Boundary[Boundary > 0] = 1
         Dapi_mask[Boundary == 1] = 0
         Dapi_mask[Dapi_mask > 0] = 255
    # remove small objects
    Dapi_mask_cleaned = morphology.remove_small_objects(np.array(Dapi_mask,bool), min_size=min_area)
    Labeled_Dapi_mask=label(Dapi_mask_cleaned, connectivity=1)
    CodexObj["Labeled_Dapi_mask"] = Labeled_Dapi_mask.astype(np.uint32)

    del DAPI_im
    del Labeled_Dapi_mask
    del Dapi_mask
    del inputpath, outputpath
    gc.enable()
    gc.collect()
    return (CodexObj)


def Membrane_seg(CodexObj, Rin=3, Rout=15):
    print('Membrane segmentation starts.....')
    startTime = time.time()
    kernel = disk(1)
    Mask_Dapi = CodexObj["Labeled_Dapi_mask"]
    M = Mask_Dapi.copy()

    for r in range(Rout):
        Mt = dilation(M, kernel)
        Mt[M != 0] = M[M != 0]
        M = Mt

    Mt = erosion(M, kernel)
    M[(M - Mt) > 0] = 0
    Membrane_mask = M
    kernel = disk(Rin)
    EM = erosion(Mask_Dapi, kernel)
    Membrane_mask[EM > 0] = 0
    CodexObj["Labeled_Membrane_mask"] = Membrane_mask.astype(np.uint32)
    print('Membrane segmentation ends.....')
    CodexObj["time_Membrane_segmentation"]= time.time() - startTime
    del Mask_Dapi, M, Mt, Membrane_mask, EM
    gc.enable()
    gc.collect()

    return (CodexObj)


def Create_cell_info_table(CodexObj):
    print('Create cell info table starts.....')
    cell_info_table = pd.DataFrame()
    cell_info_table = pd.DataFrame(columns=['cell_ID', 'sample_ID', 'ROI_ID', 'X', 'Y', 'nuclear-size', 'membrane-size', 'umap_x', 'umap_y', 'Cluster_reso_1', 'Cluster_reso_2', 'Cluster_reso_3', 'Cluster_reso_4', 'Cluster_reso_5', 'Cluster_reso_6'])
    nuclear_prop = regionprops_table(CodexObj["Labeled_Dapi_mask"], properties=('centroid', 'area'))
    Membrane_prop = regionprops_table(CodexObj["Labeled_Membrane_mask"], properties=('area',))
    cell_info_table['X'] = nuclear_prop['centroid-0']
    cell_info_table['Y'] = nuclear_prop['centroid-1']
    cell_info_table['nuclear-size'] = nuclear_prop['area']
    cell_info_table['membrane-size'] = Membrane_prop['area']
    cell_info_table['ROI_ID'] = 1
    CodexObj["cell_info_table"] = cell_info_table
    cell_info_table['cell_ID'] = range(CodexObj["cell_info_table"].__len__())
    CodexObj['cell_info_table']['sample_ID'] = CodexObj["sample_ID"]
    #Store cell_info_table in csv file
    #cell_info_table.to_csv(os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))+'/'+'cell_info_table.csv', index = False)
    print('Create cell info table ends.....')
    del cell_info_table, nuclear_prop, Membrane_prop
    gc.enable()
    gc.collect()
    return (CodexObj)


def Calculate_Average_Intensity(CodexObj, AVIT=0):
    print('Calculate average intensity from compartments [Nuclear or Membrane].....')
    startTime = time.time()
    
    if AVIT==0:
        Average_Intensity_table = pd.DataFrame()
        colno = CodexObj["marker_info_table"].shape
        colno = colno[0]
        protein_list = []

        for col in range(colno):
            protein_list.append(CodexObj["marker_info_table"].loc[col][0])

        Average_Intensity_table = pd.DataFrame(columns=protein_list,
                                           index=range(CodexObj["cell_info_table"].__len__()))
               
        for i in tqdm(range(colno)):
            protein_im = imread(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["qptiff_path"],
                                         CodexObj["sample_ID"] + '.qptiff'), key=[i])
            protein_name = CodexObj["marker_info_table"].loc[i][0]
            min_cutoff=CodexObj["marker_info_table"].loc[i]['min_cutoff']
            max_cutoff=CodexObj["marker_info_table"].loc[i]['max_cutoff']
            if min_cutoff<0:
                min_cutoff=0
            if max_cutoff>255:
                max_cutoff=255
            protein_im[protein_im<=min_cutoff]=min_cutoff
            protein_im[protein_im>=max_cutoff]=max_cutoff
            if CodexObj["marker_info_table"].loc[i][1] == 'nuclear':
                Protein_regions_Nuc = regionprops_table(CodexObj["Labeled_Dapi_mask"], protein_im,
                                                        properties=('mean_intensity',))
                Average_Intensity_table.loc[:, protein_name] = Protein_regions_Nuc['mean_intensity'].tolist()  # pseudo code

            else:
                Protein_regions_Mem = regionprops_table(CodexObj["Labeled_Membrane_mask"], protein_im,
                                                    properties=('mean_intensity',))
                Average_Intensity_table.loc[:, protein_name] = Protein_regions_Mem['mean_intensity'].tolist()  # pseudo code

        CodexObj["Average_Intensity_table"] = Average_Intensity_table
        Average_Intensity_table.to_csv(os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))+'/'+'Average_Intensity_table.csv', index = False)
        
    else:
        Average_Intensity_table = pd.read_csv(os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))+'/'+'Average_Intensity_table.csv')
        CodexObj["Average_Intensity_table"] = Average_Intensity_table
        
    # normalize average intensity table and add it to CODEXobj
    Average_Intensity_table_normalized=Average_Intensity_table.copy()
    for col in Average_Intensity_table:    
        Average_Intensity_table_normalized[col]=stats.zscore(Average_Intensity_table[col])
    
    CodexObj["Average_Intensity_table_normalized"] = Average_Intensity_table_normalized  # here it is the protDF normalized
    Average_Intensity_table_normalized.to_csv(os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))+'/'+'Average_Intensity_table_normalized.csv', index = False)

    CodexObj["time_Calculate_average_intensity"]= time.time() - startTime
    del Average_Intensity_table
    gc.enable()
    gc.collect()
    return (CodexObj)


def leiden(adata, resolution=1.0):

    # Adjacency graph
    adjacency = adata.obsp['connectivities']
    offsets = cudf.Series(adjacency.indptr)
    indices = cudf.Series(adjacency.indices)
    g = cugraph.Graph()
    if hasattr(g, 'add_adj_list'):
        g.add_adj_list(offsets, indices, None)
    else:
        g.from_cudf_adjlist(offsets, indices, None)
    
    # Cluster
    leiden_parts, _ = cugraph.leiden(g,resolution = resolution)
    
    # Format output
    clusters = leiden_parts.to_pandas().sort_values('vertex')[['partition']].to_numpy().ravel()
    clusters = pd.Categorical(clusters)    
    return clusters


def Run_scanpy(CodexObj, runby='sample' ,Reso=4, n_neighbors=50, min_dist=0.05, use='GPU'):
    print('Run scanpy is running on ' + use + '...')
    protDF = CodexObj["Average_Intensity_table"]
    cell_info_table = CodexObj["cell_info_table"]
    if runby=='sample':
        cols_to_drop = CodexObj["marker_info_table"]['use_for_clustering']
        markers_to_drop = CodexObj["marker_info_table"]['marker_names']
    
        #Here we are droping the column of any biomarker that is assigned as False in the use for clustering column in marker info table
        for i in range(CodexObj["marker_info_table"].__len__()):
            if cols_to_drop[i] == False:
                var = markers_to_drop[i]
                print('Biomaker column of: ' + var + ' is removed from the Average Intensity table ')
                protDF = protDF.drop(var, axis=1)
    
    #Here we make the spatDF similar to the format of spatDF of scanpy
    spatDF = cell_info_table.copy()
    #spatDF = spatDF.drop(['membrane-size'], axis=1)
    spatDF = spatDF[['X', 'Y', 'nuclear-size']]
    spatDF['ImageID'] = 0
    spatDF.rename(columns={'X': 'spatial_X', 'Y': 'spatial_Y', 'nuclear-size': 'Area'}, inplace=True)
    #Intialize anndata object
    adata = ad.AnnData(protDF, dtype=np.float32)
    adata.obs = spatDF
    #Apply normalization
    if runby=='sample':
        sc.pp.log1p(adata)  
    adataScaled = sc.pp.scale(adata, max_value=10, copy=True)
    # Apply PCA
    # print('Applying PCA to adataScaled...')
    # sc.tl.pca(adataScaled, svd_solver='arpack')
    # Create neighborhood graph
    print('Creating neighborhood graph...')
    startTime = time.time()
    if use=='GPU':
        sc.pp.neighbors(adataScaled, n_neighbors=n_neighbors, method='rapids') #n_pcs=15,
    else:
        sc.pp.neighbors(adataScaled, n_neighbors=n_neighbors)
    
    CodexObj["time_neighbors"]= time.time() - startTime
    # Find the umap
    startTime = time.time()
    print('Finding the UMAP...')
    if use=='GPU':
        sc.tl.umap(adataScaled, min_dist=min_dist, spread=0.5, method='rapids') 
    else:
        sc.tl.umap(adataScaled, min_dist=min_dist) 

    CodexObj["time_UMAP"]= time.time() - startTime

    #### Storing Umap coordinates in cell info table
    np.shape(adataScaled.obsm['X_umap'])
    umap = adataScaled.obsm['X_umap']
    cell_info_table['umap_x'] = umap[:, 0]
    cell_info_table['umap_y'] = umap[:, 1]

    # Perform Clustering at different resolutions 
    print('Perform Clustering...')
    startTime = time.time()   
    if use=='GPU':
        myset=range(1,7)
        # [myset.append(rdx) for rdx in range(1,7) if rdx != Reso]
        # myset.append(Reso)
    elif use=='CPU':
        myset=[Reso]
        #myset.append(Reso)
        
    for rdx in myset:
        # Run leiden
        if use=='GPU':
            adataScaled.obs['leiden_res'+str(rdx)] = leiden(adataScaled, resolution=rdx)
        else:
            adataScaled.obs['leiden_res'+str(rdx)] = sc.tl.leiden(adataScaled, resolution=rdx)
        # Add cluster ID to cell info table
        ClusterID = adataScaled.obs['leiden_res'+str(rdx)]
        cell_info_table['Cluster_reso_'+str(rdx)] = ClusterID.iloc[:].values
        # Create Cluster info table
        Cluster_info_table=pd.DataFrame()
        Cluster_info_table=pd.DataFrame(columns=['Cluster-ID','Number-of-cells','percentage-of-cells','cell-type'])
        Cluster_info_table['Cluster-ID']=np.unique(ClusterID.iloc[:].values)
        ClusterIDmylist=ClusterID.tolist()    
        counts=[]
        perscells=[]
        ClusterIDmylist.count(0)
        for i in range(len(np.unique(ClusterID.iloc[:].values))):
            counts.append(ClusterIDmylist.count(i))
            perscells.append((ClusterIDmylist.count(i)/len(cell_info_table['Cluster_reso_'+str(rdx)])))
        Cluster_info_table['Number-of-cells']=counts
        Cluster_info_table['percentage-of-cells']=perscells
        # Cluster_average_intensity_table
        Average_Intensity_table_normalized=CodexObj["Average_Intensity_table_normalized"].copy()
        if runby=='sample':
            Average_Intensity_table_normalized['Cluster-ID']=cell_info_table['Cluster_reso_'+str(rdx)]
        else:       
            temp = cell_info_table['Cluster_reso_'+str(rdx)]
            Average_Intensity_table_normalized['Cluster-ID']=temp.array
        Cluster_Average_Intensity_table=pd.DataFrame()
        Cluster_Average_Intensity_table=pd.DataFrame(columns=Average_Intensity_table_normalized.columns)
        Cluster_Average_Intensity_table['Cluster-ID']=np.unique(Average_Intensity_table_normalized['Cluster-ID'])
        for name,values in Average_Intensity_table_normalized.iloc[:, 0:-1].iteritems():
            sum=Average_Intensity_table_normalized.groupby(by=["Cluster-ID"])[name].sum() 
            count = Average_Intensity_table_normalized.groupby(by=["Cluster-ID"])[name].count()
            Cluster_Average_Intensity_table[name]=sum/count
        CodexObj["Cluster_info_table_reso_"+str(rdx)] = Cluster_info_table
        CodexObj["Cluster_average_intensity_table_reso_"+str(rdx)] = Cluster_Average_Intensity_table

    CodexObj["time_Clustering"]= time.time() - startTime
    print('Saving Data structures...')
    CodexObj["adataScaled"] = adataScaled
    CodexObj["protDF"] = protDF
    CodexObj["cell_info_table"] = cell_info_table
    
    cell_info_table_to_user=cell_info_table.copy()
    for rdx in range(1,7):
        if rdx != Reso:
            cell_info_table_to_user = cell_info_table_to_user.drop('Cluster_reso_'+str(rdx), axis=1)
    cell_info_table_to_user.rename(columns = {'Cluster_reso_'+str(Reso):'Cluster-ID'}, inplace = True)       
    cell_info_table_to_user.to_csv(os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))+'/'+'cell_info_table.csv', index = False) 
    CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)].to_csv(os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))+'/'+'Cluster_Average_Intensity_table.csv', index = False)
    CodexObj["Cluster_info_table_reso_"+str(Reso)].to_csv(os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))+'/'+'Cluster_info_table.csv', index = False)
    del protDF, spatDF, cell_info_table, adata, adataScaled
    gc.enable()
    gc.collect()
    return (CodexObj)
    #"Cluster_info_table_reso_"+str(Reso)

def plot_CorrMatrix(CodexObj, plotby='allmarkers'):
    print('Plotting Correlation Matrix...')
    Projects_path=CodexObj["project_path"]
    project_name=CodexObj['project_name']
    analysis_path=CodexObj['analysis_path']
    if plotby=='allmarkers':
        corrDF = CodexObj["Average_Intensity_table_normalized"].corr('pearson')
        ax_CorrMatrix = sns.clustermap(corrDF, vmin=-0.5,vmax=1, figsize=(int(len(CodexObj["Average_Intensity_table_normalized"].columns)/4),int(len(CodexObj["Average_Intensity_table_normalized"].columns)/4)),yticklabels=True, xticklabels=True, cmap='RdYlBu_r')
        if not os.path.exists(os.path.join(Projects_path, project_name, analysis_path, 'Figures')):
            os.makedirs(os.path.join(Projects_path, project_name, analysis_path, 'Figures'))
        ax_CorrMatrix.savefig(os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','CorrMatrix_all_markers.png')))
    else:
        protDF = CodexObj["protDF"]
        protDF_normalized=protDF.copy()
        for col in protDF:    
            protDF_normalized[col]=stats.zscore(protDF[col])
        corrDF = protDF_normalized.corr('pearson')
        ax_CorrMatrix = sns.clustermap(corrDF, vmin=-0.5,vmax=1, figsize=(int(len(CodexObj["protDF"].columns)/3),int(len(CodexObj["protDF"].columns)/3)),yticklabels=True, xticklabels=True, cmap='RdYlBu_r')
        #ax_CorrMatrix = sns.clustermap(corrDF, vmin=-0.5,vmax=1, figsize=(20,20),yticklabels=True, xticklabels=True, cmap='RdYlBu_r')
        if not os.path.exists(os.path.join(Projects_path, project_name, analysis_path, 'Figures')):
            os.makedirs(os.path.join(Projects_path, project_name, analysis_path, 'Figures'))
        ax_CorrMatrix.savefig(os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','CorrMatrix_use_for_clustering.png')))
    


def plot_umap(CodexObj, colorby='marker', Reso=4):
    Projects_path=CodexObj["project_path"]
    project_name=CodexObj['project_name']
    analysis_path=CodexObj['analysis_path']
    if colorby == 'marker':
        print('Plotting UMap by marker...')
        ax_markers_umap=sc.pl.umap(sc.pp.subsample(CodexObj["adataScaled"],fraction=1, copy=True), color=CodexObj["protDF"].columns, ncols=math.ceil(math.sqrt((len(CodexObj["Average_Intensity_table_normalized"].columns)))), cmap='viridis',return_fig =True)
        ax_markers_umap.savefig(os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Markers_umap.png')))
    elif colorby == 'cluster': #or 'sample' or 'celltype'
        print('Plotting UMap by cluster...')
        fig,ax=plt.subplots(1,1,figsize=(8,8))
        mytitle='Number of Cells = ' + str(len(CodexObj['Average_Intensity_table_normalized']))
        if len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)])>103:
            rand_colors=[]
            for j in range(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)])):
                rand_colors.append(["#"+''.join([random.choice('ABCDEF0123456789') for i in range(6)])])
            arr = np.array(rand_colors)
            rand_colors_set = set()
            for i in arr:
                rand_colors_set.update(set(i))
            CodexObj["adataScaled"].uns['leiden_res'+str(Reso)+'_colors']=list(rand_colors_set)
            #fig,ax=plt.subplots(1,1,figsize=(8,8))
        sc.pl.umap(CodexObj["adataScaled"], color=['leiden_res'+str(Reso)], ax=ax, size=2, title=mytitle,  legend_fontsize='x-small')
        fig.savefig((os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','cluster_umap.png'))), bbox_inches='tight')
    elif colorby == 'sample':
        print('Plotting UMap by sample...')
        fig,ax=plt.subplots(1,1,figsize=(8,8))
        mytitle='Number of Cells = ' + str(len(CodexObj['Average_Intensity_table']))
        sample_color = pd.DataFrame(CodexObj["cell_info_table"].sample_ID)
        adataScaled = CodexObj["adataScaled"]
        adataScaled.obs['sample_color']=[s for s in sample_color['sample_ID']] 
        sc.pl.umap(CodexObj["adataScaled"], color=['sample_color'], ax=ax, size=2, title=mytitle,  legend_fontsize='x-small')
        fig.savefig((os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','cluster_umap_per_sample.png'))), bbox_inches='tight')
    elif colorby == 'celltype':
        print('Plotting UMap by celltype...')
        fig,ax=plt.subplots(1,1,figsize=(8,8))
        mytitle='Number of Cells = ' + str(len(CodexObj['Average_Intensity_table'])) + ' , Number of Cell Types = ' + str(len(np.unique(CodexObj['cell_info_table']['cell_type'])))
        cell_type = pd.DataFrame(CodexObj["cell_info_table"].cell_type)
        adataScaled = CodexObj["adataScaled"]
        adataScaled.obs['cell_type']=[s for s in cell_type['cell_type']] 
        sc.pl.umap(CodexObj["adataScaled"], color=['cell_type'], ax=ax, size=2, title=mytitle,  legend_fontsize='x-small')
        fig.savefig((os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','cluster_umap_per_celltype.png'))), bbox_inches='tight')

def plot_HeatMap(CodexObj, Reso=4, plotby='cluster', useonly='use_of_clustering'):
    Projects_path=CodexObj["project_path"]
    project_name=CodexObj['project_name']
    analysis_path=CodexObj['analysis_path']
    if plotby=='cluster':
        # plot with all QC passed markers using the normalized average intensities
        print('Plotting Heat Map by cluster...')
        Cluster_average_intensity_table= CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)]
        Cluster_average_intensity_table = Cluster_average_intensity_table.drop(['Cluster-ID'],axis=1)
        # Drop marker if useonly 'use of clustering'
        if useonly=='use_of_clustering':
            cols_to_drop = CodexObj["marker_info_table"]['use_for_clustering']
            markers_to_drop = CodexObj["marker_info_table"]['marker_names']   
            #Here we are droping the column of any biomarker that is assigned as False in the use for clustering column in marker info table
            for i in range(CodexObj["marker_info_table"].__len__()):
                if cols_to_drop[i] == False:
                    var = markers_to_drop[i]
                    Cluster_average_intensity_table = Cluster_average_intensity_table.drop(var, axis=1)
    
        expand_w=int((len(Cluster_average_intensity_table.columns)/len(Cluster_average_intensity_table)))  if len(Cluster_average_intensity_table.columns)>len(Cluster_average_intensity_table) else 1
        expand_h=int((len(Cluster_average_intensity_table)/len(Cluster_average_intensity_table.columns)))  if len(Cluster_average_intensity_table)>len(Cluster_average_intensity_table.columns) else 1
        # zscore Cluster_average_intensity_table
        Cluster_average_intensity_table_normalized=Cluster_average_intensity_table.copy()
        for col in Cluster_average_intensity_table:    
            Cluster_average_intensity_table_normalized[col]=stats.zscore(Cluster_average_intensity_table[col])
        #ax_clustermap = sns.clustermap(np.transpose(Cluster_average_intensity_table_normalized), dendrogram_ratio=0.05, vmin=-1,vmax=1.5, figsize=(int(len(Cluster_average_intensity_table)/4)*expand_w,int(len(Cluster_average_intensity_table.columns)/4)*expand_h),yticklabels=True, xticklabels=True, cbar_pos=(-0.07, 0.8, 0.05, 0.18), cmap='RdYlBu_r')
        ax_clustermap = sns.clustermap(np.transpose(Cluster_average_intensity_table_normalized), dendrogram_ratio=0.05, vmin=-1,vmax=1.5, figsize=(int(len(Cluster_average_intensity_table)/4),int(len(Cluster_average_intensity_table.columns)/4)),yticklabels=True, xticklabels=True, cbar_pos=(-0.07, 0.8, 0.05, 0.18), cmap='RdYlBu_r')
        ax_clustermap.savefig(os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Heatmap_Clusters.png')))
        # Reorder Cluster_info_table to match the clusters order
        Cluster_info_table = CodexObj["Cluster_info_table_reso_"+str(Reso)]
        CodexObj["Cluster_info_table_reordered"] = Cluster_info_table.set_index('Cluster-ID').reindex(ax_clustermap.dendrogram_col.reordered_ind).reset_index()
        CodexObj["Cluster_info_table_reordered"].to_csv(os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))+'/'+'Cluster_info_table.csv', index = False) 
    elif plotby=='celltype':
        print('Plotting Heat Map by celltype...')
        celltype_Average_Intensity_table = CodexObj["celltype_Average_Intensity_table"+str(Reso)]
        celltype_Average_Intensity_table = celltype_Average_Intensity_table.drop(['Cell_type'],axis=1)
        
        # Drop marker if useonly 'use of clustering'
        if useonly=='use_of_clustering':
            cols_to_drop = CodexObj["marker_info_table"]['use_for_clustering']
            markers_to_drop = CodexObj["marker_info_table"]['marker_names']   
            #Here we are droping the column of any biomarker that is assigned as False in the use for clustering column in marker info table
            for i in range(CodexObj["marker_info_table"].__len__()):
                if cols_to_drop[i] == False:
                    var = markers_to_drop[i]
                    celltype_Average_Intensity_table = celltype_Average_Intensity_table.drop(var, axis=1)
 
        # zscore celltype_Average_Intensity_table
        celltype_Average_Intensity_table_normalized=celltype_Average_Intensity_table.copy()
        for col in celltype_Average_Intensity_table:    
            celltype_Average_Intensity_table_normalized[col]=stats.zscore(celltype_Average_Intensity_table[col])
        ax_clustermap = sns.clustermap(np.transpose(celltype_Average_Intensity_table_normalized).astype(float), dendrogram_ratio=0.05, vmin=-1,vmax=1.5, figsize=(int(len(celltype_Average_Intensity_table)),int(len(celltype_Average_Intensity_table.columns)/3)),yticklabels=True, xticklabels=True, cbar_pos=(-0.07, 0.8, 0.05, 0.18), cmap='RdYlBu_r')
        ax_clustermap.savefig(os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Heatmap_Celltype.png')))


def plot_ScatterSpatial(CodexObj, colorby='cluster', Reso=4):
    Projects_path=CodexObj["project_path"]
    project_name=CodexObj['project_name']
    analysis_path=CodexObj['analysis_path']
    if colorby == 'cluster': 
        print('Plotting Scatter spatial by cluster...')
        CodexObj["adataScaled"].obsm['X_spatial'] = CodexObj["adataScaled"].obs[['spatial_Y','spatial_X']].values
        CodexObj["adataScaled"].obsm['X_spatial'][:,1] = -1*CodexObj["adataScaled"].obsm['X_spatial'][:,1]
        ss = sc.pp.subsample(CodexObj["adataScaled"],fraction=.3, copy=True)
        figY = 9
        figX = int(figY*(CodexObj["adataScaled"].obs.spatial_Y.max()/CodexObj["adataScaled"].obs.spatial_X.max()))
        fig,ax=plt.subplots(1,1,figsize=(figX,figY))
        sc.pl.scatter(ss, basis='spatial',color=['leiden_res'+str(Reso)],size=8,ax=ax,title='',alpha=0.9, legend_fontsize = 'x-small')
        fig.savefig((os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Spatial.png'))), bbox_inches='tight')
    elif colorby == 'sample':
        print('Plotting Scatter spatial by sample...')
        CodexObj["adataScaled"].obsm['X_spatial'] = CodexObj["adataScaled"].obs[['spatial_Y','spatial_X']].values
        CodexObj["adataScaled"].obsm['X_spatial'][:,1] = -1*CodexObj["adataScaled"].obsm['X_spatial'][:,1]
        sample_color = pd.DataFrame(CodexObj["cell_info_table"].sample_ID)
        adataScaled = CodexObj["adataScaled"]
        adataScaled.obs['sample_color']=[s for s in sample_color['sample_ID']] 
        for sampleID in np.unique(sample_color):
            adataScaled_sampleID=adataScaled[adataScaled.obs['sample_color']==sampleID,:]
            #adataScaled_sampleID_color=adataScaled.obs['leiden_res4'][adataScaled.obs['sample_color']==sampleID]
            ss = sc.pp.subsample(adataScaled_sampleID,fraction=.3, copy=True)
            figY = 9
            figX = int(figY*(adataScaled_sampleID.obs.spatial_Y.max()/adataScaled_sampleID.obs.spatial_X.max()))
            fig,ax=plt.subplots(1,1,figsize=(figX,figY))
            sc.pl.scatter(ss, basis='spatial',color=['leiden_res'+str(Reso)],size=8,ax=ax,title='',alpha=0.9, legend_fontsize = 'x-small')
            fig.savefig((os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures',sampleID+'_Combined_Spatial.png'))), bbox_inches='tight')
    elif colorby == 'celltype':
        print('Plotting Scatter spatial by celltype...')       
        CodexObj["adataScaled"].obsm['X_spatial'] = CodexObj["adataScaled"].obs[['spatial_Y','spatial_X']].values
        CodexObj["adataScaled"].obsm['X_spatial'][:,1] = -1*CodexObj["adataScaled"].obsm['X_spatial'][:,1]
        cell_type = pd.DataFrame(CodexObj["cell_info_table"].cell_type)
        adataScaled = CodexObj["adataScaled"]
        adataScaled.obs['cell_type']=[s for s in cell_type['cell_type']] 
        ss = sc.pp.subsample(adataScaled.obs['cell_type'],fraction=.3, copy=True)
        figY = 9
        figX = int(figY*(CodexObj["adataScaled"].obs.spatial_Y.max()/CodexObj["adataScaled"].obs.spatial_X.max()))
        fig,ax=plt.subplots(1,1,figsize=(figX,figY))
        sc.pl.scatter(adataScaled, basis='spatial',color=['cell_type'],size=8,ax=ax,title='',alpha=0.9, legend_fontsize = 'x-small')
        fig.savefig((os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Celltype_Spatial.png'))), bbox_inches='tight')


def plot_Pie_chart(CodexObj, colorby='cluster', Reso=4):
    Projects_path=CodexObj["project_path"]
    project_name=CodexObj['project_name']
    analysis_path=CodexObj['analysis_path']
    if colorby == 'cluster': 
        if np.max(CodexObj["Cluster_info_table_reso_"+str(Reso)]['Cluster-ID'])<=20:
            ncol=1
            fonts=8
        elif np.max(CodexObj["Cluster_info_table_reso_"+str(Reso)]['Cluster-ID'])<=41:
            ncol=2
            fonts=11
        elif np.max(CodexObj["Cluster_info_table_reso_"+str(Reso)]['Cluster-ID'])<=62:
            ncol=3
            fonts=14
        elif np.max(CodexObj["Cluster_info_table_reso_"+str(Reso)]['Cluster-ID'])<=83:
            ncol=4
            fonts=17
        else:
            ncol=5
            fonts=20
        print('Plotting Pie chart by cluster...')
        fig,ax=plt.subplots(1,1,figsize=(ncol*7,ncol*7))
        mylabels= CodexObj["Cluster_info_table_reso_"+str(Reso)]['percentage-of-cells'] 
        #labels = [str(round(round(s,6)*100,2)) +'%' for s in CodexObj["Cluster_info_table"]['percentage-of-cells']]
        labels = [str(s) for s in CodexObj["Cluster_info_table_reso_"+str(Reso)]['Cluster-ID']]
        #CodexObj["adataScaled"].obs.leiden_res4.value_counts()
        ax.pie(CodexObj["Cluster_info_table_reso_"+str(Reso)]['Number-of-cells'],explode=[0.05]*len(CodexObj["Cluster_info_table_reso_"+str(Reso)]['Cluster-ID']),pctdistance=-0.85,startangle=0,
            labels=labels, colors=CodexObj["adataScaled"].uns['leiden_res'+str(Reso)+'_colors'])
        plt.set_cmap('Set2')
        params = {'legend.fontsize': 'small',
                 'axes.labelsize':14,
                'axes.titlesize':'small',
                'xtick.labelsize':'small',
                'ytick.labelsize':'small'}
        pylab.rcParams.update(params)
        pylab.rcParams['font.size'] = 16
        labels = [(str(s) + '-->' + str(round(round(p,6)*100,2)) +'%')  for s, p in zip(CodexObj["Cluster_info_table_reso_"+str(Reso)]['Cluster-ID'],CodexObj["Cluster_info_table_reso_"+str(Reso)]['percentage-of-cells'])]
        #plt.legend( loc = (1.1, 0.2), labels=round(round(mylabels,3)*100,2), ncol=4)
        plt.legend( loc = (1.1, 0.2), labels=labels, fontsize=fonts, ncol=ncol)
        fig = plt.gcf()
        plt.tight_layout()
        mypath=os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))
        plt.savefig(os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Pie_chart.png')),dpi=300)
    if colorby == 'sample': 
        print('Plotting Pie chart by sample...')
        fig,ax=plt.subplots(1,1,figsize=(8,8))
        sample_color = pd.DataFrame(CodexObj["cell_info_table"].sample_ID)
        labels = [str(round((s/np.sum(CodexObj["sample_nocells"]))*100,2)) +'%' for s in CodexObj["sample_nocells"]]
        ax.pie(CodexObj["sample_nocells"],explode=[0.05]*len(CodexObj["sample_nocells"]),pctdistance=-0.85,startangle=0, labels=labels)
        plt.set_cmap('Set2')
        params = {'legend.fontsize': 'small',
             'axes.labelsize':12,
            'axes.titlesize':'small',
             'xtick.labelsize':'small',
             'ytick.labelsize':'small'}
        pylab.rcParams.update(params)
        sample_IDs=CodexObj["sample_IDs"]
        labels = [str(s) for s in sample_IDs]
        plt.legend( loc = (0.95, 0.6), labels=labels, fontsize=10, ncol=1)
        fig = plt.gcf()
        plt.tight_layout()
        mypath=os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))
        plt.savefig(os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Pie_chart_per_sample.png')),dpi=300)
    if colorby == 'celltype':
        print('Plotting Pie chart by celltype...')
        fig,ax=plt.subplots(1,1,figsize=(11,11))
        labels = [str(s) for s in CodexObj["celltype_info_table_reso_"+str(Reso)]['cell_type']]
        ax.pie(CodexObj["celltype_info_table_reso_"+str(Reso)]['Number_of_cells'],explode=[0.05]*len(CodexObj["celltype_info_table_reso_"+str(Reso)]['cell_type']),pctdistance=-0.85,startangle=0,
        labels=labels, colors=CodexObj["adataScaled"].uns['cell_type_colors'])
        plt.set_cmap('Set2')
        params = {'legend.fontsize': 'small',
                 'axes.labelsize':8,
                'axes.titlesize':'small',
                'xtick.labelsize':'small',
                'ytick.labelsize':'small'}
        pylab.rcParams.update(params)
        pylab.rcParams['font.size'] = 10
        labels = [(str(s) + '-->' + str(round(round(p,6)*100,2)) +'%')  for s, p in zip(CodexObj["celltype_info_table_reso_"+str(Reso)]['cell_type'],CodexObj["celltype_info_table_reso_"+str(Reso)]['percentage_of_cells'])]
        plt.legend( loc = (1.1, 0.3), labels=labels, fontsize=9, ncol=1)
        fig = plt.gcf()
        plt.tight_layout()
        mypath=os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'Output-data'))
        plt.savefig(os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Pie_chart_Celltype.png')),dpi=300)




def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]

    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size

    # Make sure the placeholder doesn't zoom in
    placeholder.height = height 
    placeholder.width = width 

    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)

    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side
        

def Create_PPT_slides(CodexObj, runby='sample', Reso=4): 
    prs = Presentation()
    Projects_path=CodexObj["project_path"]
    project_name=CodexObj['project_name']
    analysis_path=CodexObj['analysis_path']
    if runby=='sample' or runby=='celltype':
        sample_ID=CodexObj["sample_ID"]
        print('Creating the PPT slides for one sample (patient)...')

    elif runby=='combined':
        sample_IDs=CodexObj["sample_IDs"]
        print('Creating the PPT slides for Dataset combined...')  
    # Slide 0 
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    
    placeholder = slide.placeholders[1]
    _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','CorrMatrix_all_markers.png')))
    title=slide.shapes.title.text="Correlation Matrix (all BioMarkers)"
    if runby=='sample' or runby=='celltype':
            sub=slide.placeholders[2].text="for sample " + sample_ID +  "  with " + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)].columns)-1) + " bio-marker"
    elif runby=='combined':
            sub=slide.placeholders[2].text="for combined dataset"  + " with " + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)].columns)-1) + " common bio-markers across all samples" 

    # Slide 1
    if runby=='sample' or runby=='celltype':
        slide = prs.slides.add_slide(prs.slide_layouts[8])
        placeholder = slide.placeholders[1]
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','CorrMatrix_use_for_clustering.png')))
        title=slide.shapes.title.text="Correlation Matrix (subset of BioMarkers)"   
        sub=slide.placeholders[2].text="for sample " + sample_ID +  "  with " + str(len(CodexObj["protDF"].columns)) + " bio-marker"
    
    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    placeholder = slide.placeholders[1]
    _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Markers_umap.png')))
    title=slide.shapes.title.text="UMAP for each biomarker"
    if runby=='sample' or runby=='celltype':
        sub=slide.placeholders[2].text="for sample " + sample_ID + " with " + str(len(CodexObj["protDF"].columns)) + " bio-marker"
    elif runby=='combined':
        sub=slide.placeholders[2].text="for combined dataset"  + " with " + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)].columns)-1) + " common bio-markers"

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    placeholder = slide.placeholders[1]
    #_add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Heatmap_Clusters.png')))
    #title=slide.shapes.title.text="UMap for all the clusters"
    if runby=='sample':
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','cluster_umap.png')))
        title=slide.shapes.title.text="UMAP for all the clusters"
        sub=slide.placeholders[2].text="for sample " + sample_ID + " with " + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)])-1) + " clusters"
    elif runby=='combined':
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','cluster_umap_per_sample.png')))
        title=slide.shapes.title.text="UMAP per sample"
        sub=slide.placeholders[2].text="for combined dataset"  

    elif runby=='celltype':
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','cluster_umap_per_celltype.png')))
        title=slide.shapes.title.text="UMAP for all  cell types"
        sub=slide.placeholders[2].text="for sample " + sample_ID + " with " + str(len(CodexObj["celltype_Average_Intensity_table"+str(Reso)])) + " cell types"

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    placeholder = slide.placeholders[1]
    # _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Heatmap_Clusters.png')))
    # title=slide.shapes.title.text="Correlation Matrix between BioMarkers and Clusters"
    if runby=='sample':
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Heatmap_Clusters.png')))
        title=slide.shapes.title.text="Heatmap Cluster Average Expression"
        sub=slide.placeholders[2].text="for sample " + sample_ID + " with " + str(len(CodexObj["protDF"].columns)) + " bio-marker and " + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)])-1) + " clusters" 

    elif runby=='combined':
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Heatmap_Clusters.png')))
        title=slide.shapes.title.text="Heatmap Cluster Average Expression"
        sub=slide.placeholders[2].text="for combined dataset with " + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)].columns)-1) + " common bio-marker and " + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)])-1) + " clusters" 

    elif runby=='celltype':
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Heatmap_Celltype.png')))
        title=slide.shapes.title.text="Heatmap Cluster Average Expression"
        sub=slide.placeholders[2].text="for sample " + sample_ID + " with " + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)].columns)-1) + " bio-marker and " +  str(len(CodexObj["celltype_Average_Intensity_table"+str(Reso)])) + " cell types"


    # Slide 5
    if runby=='sample':
        slide = prs.slides.add_slide(prs.slide_layouts[8])
        placeholder = slide.placeholders[1]
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Spatial.png')))
        title=slide.shapes.title.text="Single Cell Cluster Spatial Map Analysis"
        sub=slide.placeholders[2].text='With Number of Cells = ' + str(len(CodexObj["Average_Intensity_table"])) + ' and ' + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)])-1) +  " clusters"
    elif runby=='combined':
        c=0
        for sample_ID in sample_IDs:
            slide = prs.slides.add_slide(prs.slide_layouts[8])
            placeholder = slide.placeholders[1]
            _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures', sample_ID +'_Combined_Spatial.png')))
            title=slide.shapes.title.text="Single Cell Cluster Spatial Map Analysis for " + sample_ID
            sub=slide.placeholders[2].text='With Number of Cells = ' + str(CodexObj["sample_nocells"][c]) + ' and ' + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)])-1) +  " clusters"
            c=c+1
    elif runby=='celltype':
        slide = prs.slides.add_slide(prs.slide_layouts[8])
        placeholder = slide.placeholders[1]
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Celltype_Spatial.png')))
        title=slide.shapes.title.text="Single Cell type Spatial Map Analysis"
        sub=slide.placeholders[2].text='With Number of Cells = ' + str(len(CodexObj["Average_Intensity_table"])) + ' and ' + str(len(CodexObj["celltype_Average_Intensity_table"+str(Reso)])) +  " celltype"
    
    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    placeholder = slide.placeholders[1]
    if runby=='combined' or runby=='sample':
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Pie_chart.png')))
        title=slide.shapes.title.text="Pie-chart of cell per cluster percentage"
        sub=slide.placeholders[2].text='With Number of Cells = ' + str(len(CodexObj['Average_Intensity_table'])) + ' and ' + str(len(CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)])-1) +  " clusters"
    
    if runby=='combined':
        slide = prs.slides.add_slide(prs.slide_layouts[8])
        placeholder = slide.placeholders[1]
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Pie_chart_per_sample.png')))
        title=slide.shapes.title.text="Pie-chart of cell per sample percentage"
        sub=slide.placeholders[2].text='With Number of Cells = ' + str(len(CodexObj['Average_Intensity_table'])) 
    if runby=='celltype':
        _add_image(slide,1,os.path.join(os.path.join(Projects_path, project_name, analysis_path,'Figures','Pie_chart_Celltype.png')))
        title=slide.shapes.title.text="Pie-chart of cell per celltype percentage"
        sub=slide.placeholders[2].text='With Number of Cells = ' + str(len(CodexObj['Average_Intensity_table'])) + ' and ' + str(len(CodexObj["celltype_Average_Intensity_table"+str(Reso)])) +  " celltype"

    #Save the presentation
    if not os.path.exists(os.path.join(Projects_path, project_name, 'Analysis', 'Summary')):
        os.makedirs(os.path.join(Projects_path, project_name, 'Analysis', 'Summary'))
    if runby=='combined':
        pathtoppt=os.path.join(os.path.join(Projects_path, project_name, 'Analysis','Summary', 'DataSet_Combined_Figures.pptx'))
        prs.save(pathtoppt)
    elif runby=='sample':
        pathtoppt=os.path.join(os.path.join(Projects_path, project_name, 'Analysis','Summary', sample_ID +'_Figures.pptx'))
        prs.save(pathtoppt)
    elif runby=='celltype':
        pathtoppt=os.path.join(os.path.join(Projects_path, project_name, 'Analysis','Summary', sample_ID +'_Figures_after_phenotyping.pptx'))
        prs.save(pathtoppt)

def Run_dataset_combined(CODEXobj_master):
    Projects_path=CODEXobj_master["project_path"]
    project_name=CODEXobj_master["project_name"]
    analysis_path=CODEXobj_master["analysis_path"]
    sample_IDs=CODEXobj_master["sample_IDs"]

    if not os.path.exists(os.path.join(Projects_path, project_name, 'Analysis','Combined_dataset')):
           os.makedirs(os.path.join(Projects_path, project_name, 'Analysis','Combined_dataset'))

    if not os.path.exists(os.path.join(Projects_path, project_name, 'Analysis','Combined_dataset','Figures')):
           os.makedirs(os.path.join(Projects_path, project_name, 'Analysis','Combined_dataset','Figures'))

    if not os.path.exists(os.path.join(Projects_path, project_name, 'Analysis','Combined_dataset','Output-data')):
           os.makedirs(os.path.join(Projects_path, project_name, 'Analysis','Combined_dataset','Output-data'))

    # Loop over cell info table and average intensity table (normalized and has all markers)
    allframes=[]
    allcellsinfo=[]
    sample_nocells=[]
    for sample_ID in sample_IDs:
        analysis_path = os.path.join('Analysis',sample_ID)
        data_path=os.path.join('Image-data',sample_ID)
        df = pd.read_csv (os.path.join(Projects_path, project_name, analysis_path, 'Output-data', 'Average_Intensity_table_normalized.csv'))
        MInfoT= pd.read_csv (os.path.join(Projects_path, project_name, data_path, sample_ID + '_marker_info_table.csv'))
        # Consider only use for clustering
        cols_to_drop = MInfoT['use_for_clustering']
        markers_to_drop = MInfoT['marker_names']
        # Here we are droping the column of any biomarker that is assigned as False in the use for clustering column in marker info table
        for i in range(MInfoT.__len__()):
            if cols_to_drop[i] == False:
                var = markers_to_drop[i]
                df = df.drop(var, axis=1)
        
        df.columns = df.columns.str.strip().str.lower() 
        sample_nocells.append(np.shape(df)[0])
        allframes.append(df)
        df_cellinfo = pd.read_csv (os.path.join(Projects_path, project_name, analysis_path, 'Output-data', 'cell_info_table.csv'))
        df_cellinfo = df_cellinfo.drop ('cell_ID',axis=1) 
        df_cellinfo = df_cellinfo.drop ('umap_x',axis=1) 
        df_cellinfo = df_cellinfo.drop ('umap_y',axis=1) 
        df_cellinfo = df_cellinfo.drop ('Cluster-ID',axis=1) 
        allcellsinfo.append(df_cellinfo)

    # find common markers across all samples that are use_for_clustering==Yes (upper or lower case doesn't matter)
    common_cols = list(set.intersection(*(set(df.columns) for df in allframes)))
    Average_Intensity_table_master=pd.concat([df[common_cols] for df in allframes], ignore_index=True)
    cell_info_table=pd.concat(allcellsinfo, axis=0)
    # 
    cell_info_table['cell_ID'] = range(cell_info_table.__len__())
    print("Number of cells is equal to " + str(np.max(cell_info_table['cell_ID'])))
    CODEXobj_master["sample_nocells"]=sample_nocells
    CODEXobj_master["cell_info_table"]=cell_info_table
    CODEXobj_master["Average_Intensity_table"]=Average_Intensity_table_master
    Average_Intensity_table_normalized=Average_Intensity_table_master.copy()  # if you use the normalized ones, then you need to remove this normalization
    for col in Average_Intensity_table_master:
        Average_Intensity_table_normalized[col]=stats.zscore(Average_Intensity_table_master[col])
    CODEXobj_master["Average_Intensity_table_normalized"]=Average_Intensity_table_normalized

    return (CODEXobj_master)


def Save_CodexObj_as_Mat(CodexObj, runby='sample'):

    mypath=os.path.join(os.path.join(CodexObj["project_path"], CodexObj["project_name"], CodexObj["analysis_path"],'PhenoRapid_to_MIA_data'))  
    if not os.path.exists(mypath):
        os.makedirs(mypath)
    
    if runby=='sample':
        print('Saving the labeled masks for MIA .....')
        objectpath=mypath+'/'+'Labeled_Dapi_mask_'+ CodexObj["sample_ID"] + '.mat'
        CodexObj_mat = {}
        CodexObj_mat['Labeled_Dapi_mask']=CodexObj['Labeled_Dapi_mask']
        hdf5storage.savemat(objectpath, CodexObj_mat, format='7.3', matlab_compatible = True)
        objectpath=mypath+'/'+'Labeled_Membrane_mask_'+ CodexObj["sample_ID"] + '.mat'
        CodexObj_mat = {}
        CodexObj_mat['Labeled_Membrane_mask']=CodexObj["Labeled_Membrane_mask"]
        hdf5storage.savemat(objectpath, CodexObj_mat, format='7.3', matlab_compatible = True)

    print('Saving CODEX Object to mat file ...')
    CodexObj_mat = {}
    if runby=='sample':
        objectpath=mypath+'/'+'CodexObject_'+ CodexObj["sample_ID"] + '.mat'
        CodexObj_mat['sample_ID']=CodexObj["sample_ID"]
        # Marker_info_table
        CodexObj_mat['marker_info_table']=CodexObj['marker_info_table'].values
            # Labeled masks
        # CodexObj_mat['Labeled_Dapi_mask']=CodexObj['Labeled_Dapi_mask']
        # CodexObj_mat['Labeled_Membrane_mask']=CodexObj['Labeled_Membrane_mask']
        CodexObj_mat['time_Stardist_segmentation']=CodexObj['time_Stardist_segmentation']
        CodexObj_mat['time_Membrane_segmentation']=CodexObj['time_Membrane_segmentation']
        CodexObj_mat['time_Calculate_average_intensity']=CodexObj['time_Calculate_average_intensity']
        CodexObj_mat['time_sample']=CodexObj['time_sample']

    else:
        objectpath=mypath+'/'+'CodexObject_combined.mat'
        CodexObj_mat['sample_ID']=CodexObj["sample_IDs"]
        CodexObj_mat['time_sample'] = CodexObj["time_Combined"] 
   
    CodexObj_mat['project_path']=CodexObj["project_path"]
    CodexObj_mat['project_name']=CodexObj["project_name"]
    CodexObj_mat['analysis_path']=CodexObj["analysis_path"]
    # Cell_info_table
    CodexObj_mat['cell_info_table_cell_ID']=np.array(CodexObj['cell_info_table']['cell_ID'])
    # CodexObj_mat['cell_info_table_sample_ID']=CodexObj['cell_info_table']['sample_ID'].values # That works but it takes longer time
    CodexObj_mat['cell_info_table_sample_ID']=[str(namerow) for namerow in CodexObj['cell_info_table']['sample_ID'].values] # This is faster, you can reach it in matlab by ---> cell_info_table_sample_ID(i,:) 
    CodexObj_mat['cell_info_table_ROI_ID']=np.array(CodexObj['cell_info_table']['ROI_ID'])   
    CodexObj_mat['cell_info_table_X']=np.array(CodexObj['cell_info_table']['X'])
    CodexObj_mat['cell_info_table_Y']=np.array(CodexObj['cell_info_table']['Y'])
    CodexObj_mat['cell_info_table_nuclear_size']=np.array(CodexObj['cell_info_table']['nuclear-size'])
    CodexObj_mat['cell_info_table_membrane_size']=np.array(CodexObj['cell_info_table']['membrane-size'])
    CodexObj_mat['cell_info_table_umap_x']=np.array(CodexObj['cell_info_table']['umap_x'])
    CodexObj_mat['cell_info_table_umap_y']=np.array(CodexObj['cell_info_table']['umap_y'])
    CodexObj_mat['cell_info_table_Cluster_reso_1']=np.array(CodexObj['cell_info_table']['Cluster_reso_1'])
    CodexObj_mat['cell_info_table_Cluster_reso_2']=np.array(CodexObj['cell_info_table']['Cluster_reso_2'])
    CodexObj_mat['cell_info_table_Cluster_reso_3']=np.array(CodexObj['cell_info_table']['Cluster_reso_3'])
    CodexObj_mat['cell_info_table_Cluster_reso_4']=np.array(CodexObj['cell_info_table']['Cluster_reso_4'])
    CodexObj_mat['cell_info_table_Cluster_reso_5']=np.array(CodexObj['cell_info_table']['Cluster_reso_5'])
    CodexObj_mat['cell_info_table_Cluster_reso_6']=np.array(CodexObj['cell_info_table']['Cluster_reso_6'])

    # Average Intensity Table
    CodexObj_mat['Average_Intensity_table']=CodexObj['Average_Intensity_table'].values
    CodexObj_mat['Average_Intensity_table_header']=list(CodexObj['Average_Intensity_table'].columns)
    CodexObj_mat['Average_Intensity_table_normalized']=CodexObj['Average_Intensity_table_normalized'].values

    # Cluster average intensity
    CodexObj_mat['Cluster_average_intensity_table_reso_1']=CodexObj['Cluster_average_intensity_table_reso_1'].values
    CodexObj_mat['Cluster_average_intensity_table_reso_2']=CodexObj['Cluster_average_intensity_table_reso_2'].values
    CodexObj_mat['Cluster_average_intensity_table_reso_3']=CodexObj['Cluster_average_intensity_table_reso_3'].values
    CodexObj_mat['Cluster_average_intensity_table_reso_4']=CodexObj['Cluster_average_intensity_table_reso_4'].values
    CodexObj_mat['Cluster_average_intensity_table_reso_5']=CodexObj['Cluster_average_intensity_table_reso_5'].values
    CodexObj_mat['Cluster_average_intensity_table_reso_6']=CodexObj['Cluster_average_intensity_table_reso_6'].values
    # Cluster info table
    CodexObj_mat['Cluster_info_table_reso_1']=CodexObj['Cluster_info_table_reso_1'].values
    CodexObj_mat['Cluster_info_table_reso_2']=CodexObj['Cluster_info_table_reso_2'].values
    CodexObj_mat['Cluster_info_table_reso_3']=CodexObj['Cluster_info_table_reso_3'].values
    CodexObj_mat['Cluster_info_table_reso_4']=CodexObj['Cluster_info_table_reso_4'].values
    CodexObj_mat['Cluster_info_table_reso_5']=CodexObj['Cluster_info_table_reso_5'].values
    CodexObj_mat['Cluster_info_table_reso_6']=CodexObj['Cluster_info_table_reso_6'].values
    #
    CodexObj_mat['time_neighbors']=CodexObj['time_neighbors']
    CodexObj_mat['time_UMAP']=CodexObj['time_UMAP']
    CodexObj_mat['time_Clustering']=CodexObj['time_Clustering']
    scipy.io.savemat(objectpath, CodexObj_mat)    

def prepare_pheno_objects(CodexObj, Reso=4):

    analysis_path = os.path.join('Analysis',CodexObj['sample_ID'])
    Cluster_info_table_filled = pd.read_csv (os.path.join(CodexObj['project_path'], CodexObj['project_name'], analysis_path , 'Output-data', 'Cluster_info_table_filled.csv'))
    rowno = 0 
    # Create a new column in the cell info table, called, cell_type, to handle the cell type for each cell to plot Umap per celltype
    CodexObj['cell_info_table']['cell_type']=CodexObj['cell_info_table']["Cluster_reso_"+str(Reso)]
    for clusterno in Cluster_info_table_filled['Cluster-ID'].iloc[:]:
        condition = CodexObj['cell_info_table']['cell_type'].astype(str)== str(clusterno)
        CodexObj['cell_info_table']['cell_type'] = np.where(condition, Cluster_info_table_filled['cell-type'][rowno], CodexObj['cell_info_table']['cell_type'])
        rowno = rowno + 1

    # Create celltype average intensity table to plot heatmap by celltype
    rowno = 0 
    Cluster_average_intensity_table = CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)]
    for clusterno in Cluster_info_table_filled['Cluster-ID'].iloc[:]:
        condition = Cluster_average_intensity_table['Cluster-ID'].astype(str)== str(clusterno)
        Cluster_average_intensity_table['Cluster-ID'] = np.where(condition, Cluster_info_table_filled['cell-type'][rowno], Cluster_average_intensity_table['Cluster-ID'])
        rowno = rowno + 1        
    
    Cluster_average_intensity_table.rename(columns={'Cluster-ID': 'cell_type'}, inplace=True)
    celltype_Average_Intensity_table=pd.DataFrame()
    celltype_Average_Intensity_table=pd.DataFrame(columns=CodexObj["Cluster_average_intensity_table_reso_"+str(Reso)].columns[0:-1])
    celltype_Average_Intensity_table['Cell_type']=np.unique(Cluster_average_intensity_table['cell_type'])

    for name,values in celltype_Average_Intensity_table.iloc[:, 0:-1].iteritems():
        sum = Cluster_average_intensity_table.groupby(by=["cell_type"])[name].sum() 
        count = Cluster_average_intensity_table.groupby(by=["cell_type"])[name].count()
        celltype_Average_Intensity_table[name][:] = sum/count
    CodexObj["celltype_Average_Intensity_table"+str(Reso)] = celltype_Average_Intensity_table   

    # Create celltype_info_table to plot pychart
    cell_typeID = CodexObj['cell_info_table']['cell_type']
    celltype_info_table=pd.DataFrame()
    celltype_info_table=pd.DataFrame(columns=['cell_type', 'Number_of_cells','percentage_of_cells'])
    celltype_info_table['cell_type']=np.unique(cell_typeID.iloc[:].values)
    cell_typeIDmylist=cell_typeID.tolist()    
    counts=[]
    perscells=[]
    cell_typeIDmylist.count(0)
    for celltype in np.unique(cell_typeID.iloc[:].values):
        counts.append(cell_typeIDmylist.count(celltype))
        perscells.append((cell_typeIDmylist.count(celltype)/len(cell_typeID)))
    celltype_info_table['Number_of_cells']=counts
    celltype_info_table['percentage_of_cells']=perscells
    CodexObj["celltype_info_table_reso_"+str(Reso)]=celltype_info_table